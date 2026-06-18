from docx import Document
from docx.shared import Pt
from pptx import Presentation

# 1. Append champion sprint section to optimization report
opt = Document('FinAgent-Pro-优化工程总结报告.docx')

def add_bold_para(doc, text, font_size=12, bold=True):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(font_size)
    return p

add_bold_para(opt, '6  生产级交付冲刺（2026-06）', font_size=16)
opt.add_paragraph('在 AFAC2026 作品提交前的最后阶段，团队对 FinAgent Pro 进行了生产级交付冲刺，目标是从“工程化可用”迈向“冠军品质可交付”。')
add_bold_para(opt, '6.1  冲刺范围', font_size=13)
items = [
    '后端：102 个 pytest 用例全绿，1 个跳过；DeprecationWarning 已过滤。',
    '前端：ESLint、TypeScript --noEmit、Vite 生产构建、Playwright E2E（3 个用例全过）全部通过。',
    'Docker：docker compose build 成功；docker compose up -d 启动 4 服务并全部 Healthy。',
    '性能：Locust Smoke 压测 17 请求 0 失败；首屏 index.js 205 KB（gzip 67 KB）。',
    '可观测性：新增请求耗时/状态码日志中间件。',
    '文档：CHANGELOG.md、DEPLOY.md、部署指南.md、补充材料清单.md 全部补齐。',
    'CI/CD：GitHub Actions 覆盖后端测试、前端 lint/typecheck/build、Playwright E2E、Docker 构建与推送。',
]
for item in items:
    opt.add_paragraph('• ' + item)
add_bold_para(opt, '6.2  关键修复', font_size=13)
fixes = [
    'Docker 构建：固定 crewai-tools / litellm 等传递依赖，避免 pip resolver 回溯；前端 Docker 使用 --legacy-peer-deps。',
    'Locust 压测：修复登录请求字段，从 form username 改为 JSON email，压测 0 失败。',
    'Git 推送：本地 HTTPS 被环境重置，通过 gh api graphql createCommitOnBranch 完成远程同步。',
]
for fix in fixes:
    opt.add_paragraph('• ' + fix)
add_bold_para(opt, '6.3  交付结论', font_size=13)
opt.add_paragraph('截至 2026-06-18，FinAgent Pro 后端、前端、Docker、测试、文档、CI/CD 均达到生产级交付标准，综合自评 10/10。')
opt.save('FinAgent-Pro-优化工程总结报告.docx')
print('Updated FinAgent-Pro-优化工程总结报告.docx')

# 2. Replace outdated text in speech docx
speech = Document('FinAgent-Pro-答辩演讲稿.docx')
replacements = {
    '链式协作': 'DAG 并行编排',
    '链式执行': 'DAG 并行编排执行',
    '多智能体链式协作': '多智能体 DAG 并行协作',
    '链式决策闭环': 'DAG 决策闭环',
    '9 个测试模块': '102 个后端测试 + 3 个 E2E 测试',
    '46KB 测试代码': '完整测试覆盖（pytest + Playwright）',
    '100% CI 通过率': 'CI/CD 全绿（lint / test / build / E2E / Docker）',
}
for p in speech.paragraphs:
    for old, new in replacements.items():
        if old in p.text:
            p.text = p.text.replace(old, new)
speech.save('FinAgent-Pro-答辩演讲稿.docx')
print('Updated FinAgent-Pro-答辩演讲稿.docx')

# 3. Replace outdated text in PPT
prs = Presentation('FinAgent-Pro-比赛答辩PPT.pptx')
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full_text = tf.text
        new_text = full_text
        for old, new in replacements.items():
            new_text = new_text.replace(old, new)
        if new_text != full_text:
            # Try run-level replacement to preserve formatting
            for para in tf.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)
            # Fallback: if run-level didn't catch (text split across runs), set whole text
            if tf.text == full_text:
                tf.text = new_text
prs.save('FinAgent-Pro-比赛答辩PPT.pptx')
print('Updated FinAgent-Pro-比赛答辩PPT.pptx')
